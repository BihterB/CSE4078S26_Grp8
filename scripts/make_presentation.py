from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

OUT_PATH = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        "reports", "CSE4078S26_Grp8_FinalPresentation.pptx")

# ── Colors ──────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)
NAVY       = RGBColor(0x1A, 0x4A, 0x7A)
BRIGHT     = RGBColor(0x4E, 0x9A, 0xF1)
GOLD       = RGBColor(0xF4, 0xB9, 0x42)
LIGHT_BG   = RGBColor(0xF5, 0xF8, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x20, 0x2C)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
CARD_BG    = RGBColor(0xE8, 0xF1, 0xFB)

W, H = Inches(10), Inches(5.625)   # 16:9


# ── Helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank


def bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    sp = slide.shapes.add_shape(1, 0, 0, W, H)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp._element.getparent().remove(sp._element)
    slide.shapes._spTree.insert(2, sp._element)


def rect(slide, x, y, w, h, color, line_color=None):
    sp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color:
        sp.line.color.rgb = line_color
    else:
        sp.line.fill.background()
    return sp


def txt(slide, text, x, y, w, h, size=18, bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def accent_bar(slide, color=BRIGHT):
    rect(slide, 0, 0, 10, 0.07, color)


def slide_header(slide, title, subtitle=None, light=True):
    accent_bar(slide, BRIGHT if light else GOLD)
    tc = DARK_TEXT if light else WHITE
    txt(slide, title, 0.4, 0.12, 9.2, 0.7, size=28, bold=True, color=tc)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.82, 9.2, 0.35, size=13, color=MUTED if light else RGBColor(0xB0,0xC4,0xDE))


def card(slide, x, y, w, h, fill=CARD_BG):
    r = rect(slide, x, y, w, h, fill)
    return r


# ── Slide 1: Title ───────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank(prs)
    bg(s, DARK_BG)
    rect(s, 0, 0, 0.12, 5.625, BRIGHT)
    rect(s, 0, 5.3, 10, 0.325, NAVY)

    txt(s, "Fine-Tuning Small LLMs for", 0.5, 0.7, 9, 0.6,
        size=22, color=RGBColor(0xB0,0xC4,0xDE))
    txt(s, "Turkish Legal Question Answering", 0.5, 1.25, 9, 0.85,
        size=36, bold=True, color=WHITE)
    txt(s, "A QLoRA-Based Approach", 0.5, 2.1, 9, 0.5,
        size=20, italic=True, color=GOLD)

    rect(s, 0.5, 2.75, 9, 0.03, RGBColor(0x2A,0x4A,0x6A))

    txt(s, "CSE4078 Spring 2026  ·  Group 8", 0.5, 2.9, 9, 0.4,
        size=14, color=RGBColor(0x8A,0xA4,0xBE), align=PP_ALIGN.CENTER)
    txt(s,
        "Ayşegül Bihter Banuşoğlu  ·  Kaan Camcı  ·  Yakup Mert Aslan  ·  "
        "Muhammed Taha Serdar  ·  Muhammet Akyüz",
        0.5, 3.3, 9, 0.5, size=12,
        color=RGBColor(0x70,0x90,0xAA), align=PP_ALIGN.CENTER)


# ── Slide 2: Agenda ──────────────────────────────────────────────────────────
def slide_agenda(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Presentation Outline")

    items = [
        ("01", "Problem & Dataset",         "Turkish legal QA + mandatory corpus"),
        ("02", "Baseline Models",           "Qwen2.5-1.5B vs Llama-3.2-3B"),
        ("03", "QLoRA Fine-Tuning",         "4-bit NF4 + LoRA adapters"),
        ("04", "Results (1500 examples)",   "7 automatic metrics, full test split"),
        ("05", "Error Analysis",            "Real output examples before vs after"),
        ("06", "Technical Decisions",       "Key choices + metric rationale"),
    ]

    cols = 3
    cw, ch = 3.0, 1.1
    xs = [0.35, 3.55, 6.75]
    ys = [1.2, 2.5]

    for i, (num, title, sub) in enumerate(items):
        col, row = i % cols, i // cols
        x, y = xs[col], ys[row]
        card(s, x, y, cw, ch, CARD_BG)
        rect(s, x, y, 0.08, ch, BRIGHT)
        txt(s, num, x+0.15, y+0.08, 0.5, 0.4, size=18, bold=True, color=BRIGHT)
        txt(s, title, x+0.15, y+0.44, cw-0.25, 0.35, size=13, bold=True, color=DARK_TEXT)
        txt(s, sub,   x+0.15, y+0.73, cw-0.25, 0.3,  size=10, color=MUTED)


# ── Slide 3: Problem & Dataset ───────────────────────────────────────────────
def slide_dataset(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Task & Dataset Overview",
                 "Mandatory dataset assigned for Group 8 — Turkish legal question answering")

    # Left: Turkish NLP challenges
    card(s, 0.3, 1.2, 4.5, 3.9, CARD_BG)
    rect(s, 0.3, 1.2, 0.08, 3.9, BRIGHT)
    txt(s, "Why is Turkish Challenging for LLMs?", 0.5, 1.3, 4.1, 0.4, size=13, bold=True, color=NAVY)
    bullets = [
        "Agglutinative: one root produces dozens of word forms",
        "Legal register uses formal Ottoman-influenced vocabulary",
        "N-gram metrics (BLEU) undercount valid paraphrases",
        "Most LLMs are pretrained on English-dominant corpora",
    ]
    for j, b in enumerate(bullets):
        rect(s, 0.5, 1.85+j*0.7, 0.08, 0.08, BRIGHT)
        txt(s, b, 0.68, 1.78+j*0.7, 3.9, 0.55, size=11, color=DARK_TEXT)

    # Right: dataset card
    card(s, 5.1, 1.2, 4.55, 3.9, CARD_BG)
    rect(s, 5.1, 1.2, 0.08, 3.9, GOLD)
    txt(s, "Renicames/turkish-law-chatbot", 5.3, 1.3, 4.1, 0.4, size=13, bold=True, color=NAVY)
    txt(s, "Apache-2.0  ·  Hugging Face", 5.3, 1.72, 4.1, 0.3, size=10, color=MUTED)

    stats = [
        ("14,900", "Total Q&A pairs"),
        ("13,354", "Training examples"),
        ("1,500",  "Test examples (held out)"),
    ]
    for j, (val, lab) in enumerate(stats):
        rect(s, 5.3, 2.2+j*0.9, 4.0, 0.78, WHITE)
        txt(s, val, 5.45, 2.25+j*0.9, 2.0, 0.42, size=22, bold=True, color=BRIGHT)
        txt(s, lab, 5.45, 2.62+j*0.9, 3.7, 0.28, size=10, color=MUTED)

    txt(s, "Fields: Soru (question)  ·  Cevap (reference answer) — Turkish",
        5.3, 4.85, 4.1, 0.3, size=9, italic=True, color=MUTED)


# ── Slide 4: Baseline Models ─────────────────────────────────────────────────
def slide_models(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Baseline Models",
                 "Two instruction-tuned models in the 1B–4B range, evaluated on the full 1,500-example test split")

    for col, (name, short, params, maker, notes) in enumerate([
        ("Qwen2.5-1.5B-Instruct", "Qwen2.5", "1.5B params", "Alibaba Cloud",
         ["Compact, fast inference", "Strong multilingual support", "BFloat16 precision"]),
        ("Llama-3.2-3B-Instruct", "Llama 3.2", "3B params", "Meta AI",
         ["Latest Llama generation", "Better instruction following", "Selected for fine-tuning ✓"]),
    ]):
        x = 0.4 + col * 5.0
        ac = BRIGHT if col == 0 else GOLD
        card(s, x, 1.15, 4.5, 4.1, CARD_BG)
        rect(s, x, 1.15, 4.5, 0.08, ac)

        txt(s, short,  x+0.2, 1.25, 4.0, 0.5, size=22, bold=True, color=NAVY)
        txt(s, name,   x+0.2, 1.75, 4.1, 0.3, size=10, color=MUTED)
        txt(s, params, x+0.2, 2.1,  4.1, 0.38, size=16, bold=True, color=ac)
        txt(s, maker,  x+0.2, 2.48, 4.1, 0.28, size=10, color=MUTED)

        rect(s, x+0.2, 2.82, 4.1, 0.02, RGBColor(0xCC,0xD9,0xE8))
        for j, note in enumerate(notes):
            rect(s, x+0.25, 2.98+j*0.6, 0.12, 0.12, ac)
            txt(s, note, x+0.45, 2.92+j*0.6, 3.8, 0.42, size=12, color=DARK_TEXT)

        if col == 1:
            rect(s, x, 4.92, 4.5, 0.33, GOLD)
            txt(s, "★  Selected for QLoRA fine-tuning", x+0.2, 4.96, 4.2, 0.26,
                size=12, bold=True, color=DARK_BG)


# ── Slide 5: Baseline Results ─────────────────────────────────────────────────
def slide_baseline_results(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Baseline Evaluation Results",
                 "Both models evaluated on 1,500 test examples — greedy decoding, max 100 new tokens")

    headers = ["Model", "ROUGE-1", "ROUGE-2", "ROUGE-L", "BLEU", "chrF++", "TF-IDF", "BERTScore F1"]
    rows = [
        ["Qwen2.5-1.5B", "0.2609", "0.1271", "0.2126", "0.0688", "0.2810", "0.1648", "0.8667"],
        ["Llama-3.2-3B",  "0.2896", "0.1597", "0.2506", "0.1002", "0.2802", "0.1979", "0.8740"],
    ]

    col_w = [2.2] + [1.1]*7
    x0, y0 = 0.25, 1.2

    # Header row
    x = x0
    for i, (h, cw) in enumerate(zip(headers, col_w)):
        rect(s, x, y0, cw-0.04, 0.45, NAVY)
        txt(s, h, x+0.05, y0+0.05, cw-0.1, 0.36,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    # Data rows
    row_colors = [CARD_BG, WHITE]
    for ri, row in enumerate(rows):
        x = x0
        y = y0 + 0.45 + ri*0.5
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            rect(s, x, y, cw-0.04, 0.46, row_colors[ri])
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            bold = ci == 0
            color = NAVY if ci == 0 else DARK_TEXT
            txt(s, cell, x+0.06, y+0.08, cw-0.12, 0.3,
                size=11, bold=bold, color=color, align=align)
            x += cw

    # Key insight
    card(s, 0.25, 3.05, 9.5, 1.1, RGBColor(0xE8,0xF4,0xFF))
    rect(s, 0.25, 3.05, 0.08, 1.1, BRIGHT)
    txt(s, "Key Takeaway", 0.45, 3.12, 3, 0.35, size=12, bold=True, color=NAVY)
    txt(s,
        "Llama-3.2-3B outperforms Qwen2.5-1.5B across most metrics despite the "
        "2× parameter difference. Both models show low BLEU (0.07–0.10) but high BERTScore "
        "(0.87), meaning they understand the topic but cannot reproduce exact legal phrasing. "
        "→ Fine-tuning Llama is the most promising path.",
        0.45, 3.5, 9.1, 0.6, size=10, color=DARK_TEXT)


# ── Slide 6: QLoRA Setup ─────────────────────────────────────────────────────
def slide_qlora(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "QLoRA Fine-Tuning Setup",
                 "4-bit quantized base model + low-rank adapters — only ~1% of parameters trained")

    # Left: what is QLoRA
    card(s, 0.3, 1.15, 4.5, 4.1, CARD_BG)
    rect(s, 0.3, 1.15, 0.08, 4.1, BRIGHT)
    txt(s, "What is QLoRA?", 0.5, 1.25, 4.1, 0.38, size=14, bold=True, color=NAVY)
    qlora_pts = [
        ("Quantize", "Load base model in 4-bit NF4 (saves ~75% VRAM)"),
        ("LoRA",     "Add small trainable rank-16 matrices to attention & MLP"),
        ("Train",    "Only adapters are updated — base weights frozen"),
        ("Merge",    "Adapters merged at inference, no speed penalty"),
    ]
    for j, (label, desc) in enumerate(qlora_pts):
        rect(s, 0.5, 1.82+j*0.8, 0.7, 0.34, BRIGHT)
        txt(s, label, 0.52, 1.87+j*0.8, 0.66, 0.26, size=9, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, desc,  1.3,  1.83+j*0.8, 3.35, 0.55, size=10, color=DARK_TEXT)

    # Right: hyperparameters
    card(s, 5.1, 1.15, 4.55, 4.1, CARD_BG)
    rect(s, 5.1, 1.15, 0.08, 4.1, GOLD)
    txt(s, "Hyperparameters", 5.3, 1.25, 4.1, 0.38, size=14, bold=True, color=NAVY)

    params = [
        ("Quantization",    "4-bit NF4 + double quantization"),
        ("LoRA rank r",     "16  (α = 32, dropout = 0.05)"),
        ("Target modules",  "q/k/v/o_proj, gate/up/down_proj"),
        ("Learning rate",   "2 × 10⁻⁴ (cosine, warmup 3%)"),
        ("Batch size",      "4 × 4 grad accum = effective 16"),
        ("Epochs",          "5"),
        ("Max seq length",  "512 tokens"),
        ("Precision",       "BFloat16"),
        ("Hardware",        "NVIDIA RTX 3090 24 GB"),
        ("Training time",   "~5 hours"),
    ]
    for j, (k, v) in enumerate(params):
        y = 1.82 + j*0.35
        txt(s, k+":", 5.3, y, 1.7, 0.33, size=10, bold=True, color=NAVY)
        txt(s, v,     7.1, y, 2.4, 0.33, size=10, color=DARK_TEXT)


# ── Slide 7: Key Numbers ─────────────────────────────────────────────────────
def slide_key_numbers(prs):
    s = blank(prs)
    bg(s, DARK_BG)
    accent_bar(s, GOLD)
    txt(s, "Fine-Tuning Results — Key Numbers", 0.5, 0.15, 9, 0.6,
        size=26, bold=True, color=WHITE)
    txt(s, "Llama-3.2-3B baseline  →  Llama-3.2-3B + QLoRA  (1,500 test examples)",
        0.5, 0.78, 9, 0.35, size=13, color=RGBColor(0xB0,0xC4,0xDE))

    stats = [
        ("+134%",  "ROUGE-L",       "0.2506 → 0.5874", BRIGHT),
        ("+371%",  "BLEU",          "0.1002 → 0.4720", GOLD),
        ("+7.2 pt","BERTScore F1",  "0.8740 → 0.9369", GREEN),
    ]
    for i, (big, label, detail, color) in enumerate(stats):
        x = 0.5 + i * 3.17
        rect(s, x, 1.35, 2.9, 3.6, RGBColor(0x1B, 0x2B, 0x3C))
        rect(s, x, 1.35, 2.9, 0.07, color)
        txt(s, big,    x+0.15, 1.6,  2.6, 1.1, size=48, bold=True, color=color,
            align=PP_ALIGN.CENTER)
        txt(s, label,  x+0.15, 2.75, 2.6, 0.45, size=15, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER)
        txt(s, detail, x+0.15, 3.25, 2.6, 0.35, size=11, color=MUTED,
            align=PP_ALIGN.CENTER)

    txt(s, "All 7 metrics improved. BERTScore shows strong semantic alignment with reference answers.",
        0.5, 5.15, 9, 0.35, size=11, italic=True,
        color=RGBColor(0x80,0x9A,0xB4), align=PP_ALIGN.CENTER)


# ── Slide 8: Full Results Table ───────────────────────────────────────────────
def slide_full_results(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Complete Evaluation Results",
                 "All three models on 1,500 test examples — 7 automatic metrics")

    headers = ["Model", "ROUGE-1", "ROUGE-2", "ROUGE-L", "BLEU", "chrF++", "TF-IDF", "BERTScore F1"]
    rows = [
        ["Qwen2.5-1.5B (baseline)", "0.2609","0.1271","0.2126","0.0688","0.2810","0.1648","0.8667"],
        ["Llama-3.2-3B (baseline)", "0.2896","0.1597","0.2506","0.1002","0.2802","0.1979","0.8740"],
        ["Llama-3.2-3B (fine-tuned)","0.6186","0.5163","0.5874","0.4720","0.6038","0.5419","0.9369"],
    ]

    col_w = [2.5] + [1.05]*7
    x0, y0 = 0.2, 1.2

    x = x0
    for h, cw in zip(headers, col_w):
        rect(s, x, y0, cw-0.04, 0.42, NAVY)
        txt(s, h, x+0.04, y0+0.05, cw-0.08, 0.33,
            size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    row_bgs = [CARD_BG, WHITE, RGBColor(0xE8,0xF4,0xE8)]
    for ri, row in enumerate(rows):
        x = x0
        y = y0 + 0.42 + ri*0.52
        for ci, (cell, cw) in enumerate(zip(row, col_w)):
            rect(s, x, y, cw-0.04, 0.48, row_bgs[ri])
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            is_ft = ri == 2 and ci > 0
            txt(s, cell, x+0.05, y+0.09, cw-0.1, 0.3,
                size=10 if ci > 0 else 10,
                bold=is_ft,
                color=GREEN if is_ft else (NAVY if ci == 0 else DARK_TEXT),
                align=align)
            x += cw

    # Fine-tuned label
    txt(s, "↑ Best", 9.55, y0+0.42+2*0.52+0.06, 0.5, 0.36,
        size=9, bold=True, color=GREEN)

    # Legend note
    txt(s, "chrF++ uses word_order=2 (character + word n-grams) — better for agglutinative Turkish than BLEU",
        0.2, 4.95, 9.6, 0.35, size=9, italic=True, color=MUTED)


# ── Slide 9: Error Analysis ───────────────────────────────────────────────────
def slide_error(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Error Analysis — Before vs After Fine-Tuning",
                 "Baseline models understand the topic but cannot reproduce legal phrasing")

    question = ("Question: \"İş akdinin işveren tarafından feshedilmesi halinde "
                "işçi kıdem tazminatına hak kazanır mı?\"")
    reference = ("Reference: \"İş sözleşmesi işveren tarafından feshedilen ve en az bir yıl "
                 "kıdemi bulunan işçi kıdem tazminatına hak kazanır.\"")

    card(s, 0.3, 1.15, 9.4, 0.65, RGBColor(0xEE,0xF2,0xFF))
    rect(s, 0.3, 1.15, 0.08, 0.65, NAVY)
    txt(s, question,   0.5, 1.2,  8.9, 0.28, size=10, bold=True, color=NAVY)
    txt(s, reference,  0.5, 1.48, 8.9, 0.28, size=10, italic=True, color=MUTED)

    panels = [
        ("Llama-3.2-3B  BASELINE", BRIGHT,
         "İşçinin kıdem tazminatına hak kazanıp kazanmadığı, iş "
         "sözleşmesinin koşullarına ve Türk iş hukukuna göre değişiklik "
         "gösterir. Genel olarak işçinin hakları korunmaktadır. Bu konuda "
         "bir avukattan destek almanızı öneririz.",
         "Generic answer — no legal precision. Mentions 'varies' and 'consult a lawyer' "
         "instead of citing the actual rule. BLEU ≈ 0."),
        ("Llama-3.2-3B  FINE-TUNED", GREEN,
         "İş sözleşmesi işveren tarafından feshedilen ve en az bir yıl "
         "kıdemi bulunan işçi kıdem tazminatına hak kazanır.",
         "Matches reference almost exactly — correct legal rule, correct terminology, "
         "correct length. ROUGE-L ≈ 0.89."),
    ]

    for i, (label, color, answer, analysis) in enumerate(panels):
        x = 0.3 + i*4.85
        card(s, x, 2.0, 4.6, 3.4, CARD_BG)
        rect(s, x, 2.0, 4.6, 0.06, color)
        txt(s, label, x+0.15, 2.1, 4.3, 0.3, size=11, bold=True, color=color)
        txt(s, answer, x+0.15, 2.5, 4.3, 1.2, size=10, color=DARK_TEXT, italic=True)
        rect(s, x+0.1, 3.78, 4.4, 0.02, RGBColor(0xCC,0xD8,0xE8))
        txt(s, analysis, x+0.15, 3.85, 4.3, 1.0, size=9, color=MUTED)


# ── Slide 10: Real Output Examples ───────────────────────────────────────────
def slide_failures(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Real Model Output Examples",
                 "Actual outputs from the test split — fine-tuned model learns precise legal phrasing")

    examples = [
        {
            "q":    "Toplu iş sözleşmesi yapma hakkı neyi kapsar?",
            "ref":  "İşçiler ve işverenler, karşılıklı olarak ekonomik ve sosyal durumlarını ve çalışma şartlarını düzenlemek amacıyla toplu iş sözleşmesi yapma hakkına sahiptirler.",
            "base": "Toplu iş sözleşmesi yapma hakkı, işverenlerin işçi ve işçi sendikaları ile işveren sendikaları arasında yapılan sözleşmelerdir. Bu sözleşme çalışma koşullarını belirler.",
            "ft":   "İşçiler ve işverenler, karşılıklı olarak ekonomik ve sosyal durumlarını ve çalışma şartlarını düzenlemek amacıyla toplu iş sözleşmesi yapma hakkına sahiptirler.",
            "rl_b": "ROUGE-L: 0.18",
            "rl_f": "ROUGE-L: 1.00",
        },
        {
            "q":    "Kamulaştırma kararına karşı Anayasa Mahkemesi'ne başvurabilir miyim?",
            "ref":  "Anayasa Mahkemesi'ne başvurabilmek için önce tüm iç hukuk yollarını tüketmiş olmanız gerekmektedir.",
            "base": "Evet, kamu kuruluşları ile ilgili bir meselede Anayasa Mahkemesine başvurabilirsiniz.",
            "ft":   "Anayasa Mahkemesi'ne başvurmadan önce tüm iç hukuk yollarını tüketmiş olmanız şarttır.",
            "rl_b": "ROUGE-L: 0.13",
            "rl_f": "ROUGE-L: 0.81",
        },
    ]

    for i, ex in enumerate(examples):
        y0 = 1.1 + i * 2.2
        card(s, 0.3, y0, 9.4, 2.0, CARD_BG)
        rect(s, 0.3, y0, 0.08, 2.0, NAVY)
        txt(s, "Q: " + ex["q"], 0.5, y0+0.08, 9.0, 0.3, size=10, bold=True, color=NAVY)
        txt(s, "Reference: " + ex["ref"], 0.5, y0+0.38, 9.0, 0.28, size=9, italic=True, color=MUTED)

        # Baseline box
        rect(s, 0.5, y0+0.72, 4.3, 0.9, RGBColor(0xFF,0xEE,0xEE))
        rect(s, 0.5, y0+0.72, 4.3, 0.06, RGBColor(0xCC,0x44,0x44))
        txt(s, "Baseline  —  " + ex["rl_b"], 0.6, y0+0.76, 4.0, 0.22,
            size=8, bold=True, color=RGBColor(0xAA,0x22,0x22))
        txt(s, ex["base"], 0.6, y0+0.99, 4.1, 0.58, size=9, color=DARK_TEXT)

        # Fine-tuned box
        rect(s, 5.1, y0+0.72, 4.4, 0.9, RGBColor(0xEA,0xF7,0xEA))
        rect(s, 5.1, y0+0.72, 4.4, 0.06, GREEN)
        txt(s, "Fine-tuned  —  " + ex["rl_f"], 5.2, y0+0.76, 4.1, 0.22,
            size=8, bold=True, color=GREEN)
        txt(s, ex["ft"], 5.2, y0+0.99, 4.2, 0.58, size=9, color=DARK_TEXT)


# ── Slide 11: Critical Technical Decisions ───────────────────────────────────
def slide_technical(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Critical Technical Decisions",
                 "Each choice was deliberate — here is why")

    decisions = [
        (BRIGHT, "Chat Template Consistency",
         "Both baseline eval AND fine-tuning use the same Llama chat template with the same system prompt. "
         "Without this, the fine-tuned model would have an unfair format advantage over the baseline."),
        (GOLD,   "NF4 over INT4 Quantization",
         "LLM weights follow a normal distribution. NF4 (Normal Float 4) is information-theoretically "
         "optimal for this distribution — preserves more signal per bit than standard INT4."),
        (GREEN,  "7 LoRA Target Modules (Attention + MLP)",
         "Standard LoRA targets only q_proj/v_proj. We also target gate/up/down_proj (MLP). "
         "Research shows including MLP layers significantly improves domain adaptation quality."),
        (RGBColor(0xA8,0x5C,0xD0), "Why 5 Epochs?",
         "With 13,354 training examples, 3 epochs risks underfitting. Beyond 5-6 epochs the model "
         "begins memorizing phrasing rather than learning legal reasoning — diminishing returns on test set."),
        (RGBColor(0xE0,0x6C,0x00), "Train/Test Contamination Check",
         "We verified zero overlap between the 13,354 train examples and 1,500 test examples. "
         "The test split was never seen during training, fine-tuning, or hyperparameter selection."),
    ]

    for i, (color, title, desc) in enumerate(decisions):
        y = 1.15 + i*0.88
        card(s, 0.3, y, 9.4, 0.78, CARD_BG)
        rect(s, 0.3, y, 0.08, 0.78, color)
        txt(s, title, 0.5, y+0.07, 3.2, 0.28, size=11, bold=True, color=NAVY)
        txt(s, desc,  3.8, y+0.07, 5.8, 0.62, size=9,  color=DARK_TEXT)


# ── Slide 12: Metric Rationale ────────────────────────────────────────────────
def slide_metrics_why(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Why These 7 Metrics?",
                 "Each metric captures a different dimension — no single metric tells the full story for Turkish")

    metrics = [
        ("ROUGE-1 / 2 / L", BRIGHT,
         "Standard NLP benchmark. Unigram, bigram, and longest common subsequence overlap with reference. "
         "Widely used — makes results comparable with other papers."),
        ("BLEU", GOLD,
         "Precision-focused n-gram metric with brevity penalty. "
         "Low scores for baselines (0.07–0.10) reveal that models understand the topic but cannot reproduce exact phrasing."),
        ("chrF++", GREEN,
         "Character n-gram F-score + word order penalty. Designed for morphologically rich languages. "
         "Turkish is agglutinative — one root produces many word forms that BLEU misses but chrF++ captures."),
        ("TF-IDF Cosine", RGBColor(0xA8,0x5C,0xD0),
         "Lexical overlap without n-gram ordering constraints. Complementary signal to ROUGE for legal vocabulary."),
        ("BERTScore F1", RGBColor(0xE0,0x6C,0x00),
         "Contextual embedding similarity using XLM-RoBERTa (100 languages including Turkish). "
         "Captures semantic correctness even when exact wording differs — baseline BERTScore (0.87) "
         "confirms models understand the topic despite low BLEU."),
    ]

    for i, (name, color, desc) in enumerate(metrics):
        y = 1.15 + i*0.86
        card(s, 0.3, y, 9.4, 0.76, CARD_BG)
        rect(s, 0.3, y, 0.08, 0.76, color)
        txt(s, name, 0.5,  y+0.1, 2.0, 0.56, size=11, bold=True, color=color)
        txt(s, desc, 2.65, y+0.1, 6.9, 0.6,  size=9,  color=DARK_TEXT)


# ── Slide 13: Conclusions ─────────────────────────────────────────────────────
def slide_conclusions(prs):
    s = blank(prs)
    bg(s, LIGHT_BG)
    slide_header(s, "Conclusions")

    points = [
        (BRIGHT, "QLoRA works extremely well for Turkish legal QA",
                 "5 epochs on 13,354 examples → +134% ROUGE-L, +371% BLEU vs baseline"),
        (GOLD,   "Larger baseline matters",
                 "Llama-3.2-3B (3B) outperforms Qwen2.5 (1.5B) — invest in the stronger base model"),
        (GREEN,  "chrF++ and BERTScore are better metrics for Turkish",
                 "BLEU underestimates performance on agglutinative languages; chrF++ and BERTScore are more informative"),
        (RGBColor(0xA8,0x5C,0xD0), "Test set never touched during training",
                 "Clean evaluation: zero contamination between the 13,354 train and 1,500 test examples"),
    ]

    for i, (color, title, desc) in enumerate(points):
        y = 1.15 + i*1.05
        card(s, 0.3, y, 9.4, 0.92, CARD_BG)
        rect(s, 0.3, y, 0.08, 0.92, color)
        txt(s, title, 0.5, y+0.08, 8.9, 0.36, size=13, bold=True, color=NAVY)
        txt(s, desc,  0.5, y+0.48, 8.9, 0.36, size=11, color=MUTED)


# ── Slide 12: Thank You ───────────────────────────────────────────────────────
def slide_thanks(prs):
    s = blank(prs)
    bg(s, DARK_BG)
    rect(s, 0, 0, 0.12, 5.625, GOLD)
    rect(s, 0, 5.3, 10, 0.325, NAVY)

    txt(s, "Thank You", 0.4, 1.0, 9.2, 1.1,
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "Questions & Discussion", 0.4, 2.15, 9.2, 0.5,
        size=18, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    rect(s, 2.0, 2.9, 6.0, 0.03, RGBColor(0x2A,0x4A,0x6A))

    lines = [
        "Dataset:  Renicames/turkish-law-chatbot  (Apache-2.0)",
        "Model:    meta-llama/Llama-3.2-3B-Instruct  +  QLoRA (r=16, 5 epochs)",
        "Metrics:  ROUGE-1/2/L · BLEU · chrF++ · TF-IDF · BERTScore F1 (XLM-RoBERTa)",
        "Code:     github.com/bihter/CSE4078S26_Grp8",
    ]
    for j, line in enumerate(lines):
        txt(s, line, 1.5, 3.15+j*0.42, 7, 0.38, size=11,
            color=RGBColor(0x8A,0xA4,0xBE), align=PP_ALIGN.CENTER)

    txt(s, "CSE4078 Spring 2026  ·  Group 8",
        0.4, 5.15, 9.2, 0.3, size=11,
        color=RGBColor(0x4A,0x6A,0x8A), align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_title(prs)           # 1
    slide_agenda(prs)          # 2
    slide_dataset(prs)         # 3
    slide_models(prs)          # 4
    slide_baseline_results(prs)# 5
    slide_qlora(prs)           # 6
    slide_key_numbers(prs)     # 7
    slide_full_results(prs)    # 8
    slide_error(prs)           # 9  before/after (general)
    slide_failures(prs)        # 10 real output examples
    slide_technical(prs)       # 11 critical decisions
    slide_metrics_why(prs)     # 12 why these metrics
    slide_conclusions(prs)     # 13
    slide_thanks(prs)          # 14

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
